"""Generate clean PPT - 6 slides, minimal icons"""
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

GREEN = RGBColor(0x2E, 0x7D, 0x32)
GREEN_L = RGBColor(0x4C, 0xAF, 0x50)
GREEN_D = RGBColor(0x1B, 0x5E, 0x20)
GREEN_P = RGBColor(0xE8, 0xF5, 0xE9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x88, 0x88, 0x88)
LGRAY = RGBColor(0xF5, 0xF5, 0xF5)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def bg(s, c=WHITE):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = c

def tb(s, l, t, w, h, text, sz=18, b=False, c=DARK, a=PP_ALIGN.LEFT):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(sz); p.font.bold = b
    p.font.color.rgb = c; p.font.name = '微软雅黑'; p.alignment = a
    return tf

def tbm(s, l, t, w, h, lines, sz=16, c=DARK):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(sz); p.font.color.rgb = c
        p.font.name = '微软雅黑'; p.space_after = Pt(4)
    return tf

def bullet(s, l, t, w, h, items, sz=15, c=DARK):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f'· {item}'
        p.font.size = Pt(sz); p.font.color.rgb = c
        p.font.name = '微软雅黑'; p.space_after = Pt(6)
    return tf

def rect(s, l, t, w, h, fill, line=None):
    shape = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    if line: shape.line.color.rgb = line; shape.line.width = Pt(1)
    else: shape.line.fill.background()
    return shape

def bar(s, t, c=GREEN_L):
    rect(s, 1, t, 11.333, 0.04, c)

def circle(s, l, t, d, fill):
    shape = s.shapes.add_shape(9, Inches(l), Inches(t), Inches(d), Inches(d))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape

# ═══════════════ SLIDE 1: COVER ═══════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)
circle(s, 9.5, -1.5, 5.5, GREEN_P)
rect(s, 1.5, 3.8, 3, 0.06, GREEN_L)

tb(s, 1.5, 1.5, 10, 1.3, '拾  句', sz=64, b=True, c=GREEN)
tb(s, 1.5, 2.9, 10, 0.6, '一词入 · 千年应', sz=26, c=GREEN_L)
tb(s, 1.5, 4.2, 10, 0.5, 'AI 智能体开发 · 期末项目展示', sz=18, c=GRAY)
tb(s, 1.5, 4.9, 10, 0.5, '王  宸    2025212714', sz=16, c=DARK)
tb(s, 1.5, 5.6, 10, 0.4, '北京邮电大学 · 2025-2026学年第二学期', sz=13, c=GRAY)

# ═══════════════ SLIDE 2: 项目简介 ═══════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)
bar(s, 1.3)
tb(s, 1, 0.3, 10, 0.8, '项目简介', sz=32, b=True, c=GREEN)

# Left
rect(s, 0.8, 1.8, 5.5, 4.8, LGRAY)
tb(s, 1.2, 2.0, 5, 0.5, '做什么', sz=20, b=True, c=GREEN_D)
rect(s, 1.2, 2.55, 5, 0.02, GREEN_L)
bullet(s, 1.2, 2.9, 5, 3.5, [
    '输入一个词（孤独/暗恋/毕业...）',
    'AI 从唐诗宋词古文散文中匹配最贴切的古典名句',
    '生成一张诗意卡片，可直接截图分享',
    '不是让AI写诗，而是让AI帮你找诗',
    '李白杜甫早就写好了，你只需要找到它',
], sz=14, c=DARK)

# Right
rect(s, 7, 1.8, 5.5, 4.8, GREEN_P)
tb(s, 7.4, 2.0, 5, 0.5, '特点', sz=20, b=True, c=GREEN_D)
rect(s, 7.4, 2.55, 5, 0.02, GREEN_L)
bullet(s, 7.4, 2.9, 5, 3.5, [
    '四种文学风格：唐诗/宋词/古文/散文',
    '一词四境 —— 同一情绪，四种古典回应',
    '极简交互：3个元素，3秒上手',
    '视觉即内容：每种风格独立CSS视觉方案',
    '免费托管，永久在线',
], sz=14, c=DARK)

# ═══════════════ SLIDE 3: 技术架构 ═══════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)
bar(s, 1.3)
tb(s, 1, 0.3, 10, 0.8, '技术架构', sz=32, b=True, c=GREEN)

# Left: flow diagram
layers = [
    ('Streamlit 前端', '输入框 · 风格选择 · CSS卡片渲染', GREEN_P),
    ('Agent 层', 'generate_card() · JSON解析 · Prompt注入', RGBColor(0xC8, 0xE6, 0xC9)),
    ('LLM Client', 'OpenAI SDK封装 · 超时处理 · 异常捕获', RGBColor(0xA5, 0xD6, 0xA7)),
    ('SiliconFlow API', 'DeepSeek-V3 大语言模型', GREEN_L),
]
for i, (name, desc, color) in enumerate(layers):
    y = 1.8 + i * 1.3
    rect(s, 1, y, 5.5, 1.0, color)
    c = WHITE if i == 3 else GREEN_D
    tb(s, 1.2, y+0.1, 5, 0.4, name, sz=18, b=True, c=c)
    tb(s, 1.2, y+0.55, 5, 0.4, desc, sz=12, c=c if i == 3 else DARK)
    if i < 3:
        tb(s, 3.5, y+1.0, 0.5, 0.4, '▼', sz=14, c=GREEN_L, a=PP_ALIGN.CENTER)

# Right: tech stack
bullet(s, 7.5, 1.8, 5, 5, [
    'Python 3.x + Streamlit 纯Python Web框架',
    'DeepSeek-V3 大模型（硅基流动API调用）',
    'OpenAI SDK 兼容格式，易扩展',
    'Streamlit Cloud 免费托管，GitHub Push 自动部署',
    'Secrets 管理 API Key，代码不留敏感信息',
    '60s 超时 + try/catch 异常处理，不会白屏',
    '纯 CSS 装饰效果，无外部图片依赖',
], sz=14, c=DARK)

# ═══════════════════════════════════════
# SLIDE 4: 核心创新
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)
bar(s, 1.3)
tb(s, 1, 0.3, 10, 0.8, '核心创新', sz=32, b=True, c=GREEN)

innovations = [
    ('AI 不写诗，AI 找诗',
     '拒绝AI创作，转为语义检索。利用LLM内化的古典文学知识，连接"现代关键词"与"古典名句"。本质是"有温度的搜索"。强制JSON输出，source字段标注作者和篇名，杜绝幻觉。'),
    ('一词四境，横跨千年',
     '同一关键词在唐诗、宋词、古文、散文四种体裁中获得完全不同的美学表达。切换风格按钮，即在同一情绪中穿行四种文学传统。每种风格有独立的CSS视觉方案。'),
    ('极简交互，产品即内容',
     '全界面仅3个可操作元素。零学习成本，不写Prompt不调参数。卡片设计精美，截图即可分享。'),
    ('视觉设计是文学体验的一部分',
     '唐诗配宣纸印章，宋词配淡紫圆环，古文配纸色邮戳，散文配暖黄衬线。纯CSS实现，视觉服务于内容而非装饰。'),
]
for i, (title, desc) in enumerate(innovations):
    y = 1.8 + i * 1.4
    rect(s, 0.8, y, 11.7, 1.2, LGRAY if i % 2 == 0 else GREEN_P)
    tb(s, 1.2, y+0.1, 11, 0.4, title, sz=18, b=True, c=GREEN_D)
    tb(s, 1.2, y+0.55, 11, 0.6, desc, sz=13, c=DARK)

# ═══════════════════════════════════════
# SLIDE 5: 问题与反思
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)
bar(s, 1.3)
tb(s, 1, 0.3, 10, 0.8, '关键问题与反思', sz=32, b=True, c=GREEN)

# Left: problems
tb(s, 1, 1.8, 5.5, 0.5, '遇到的问题', sz=20, b=True, c=GREEN_D)
rect(s, 1, 2.35, 5.5, 0.02, GREEN_L)
problems = [
    ('AI 编造假诗', '四赛道Prompt限定作者 + JSON结构化 + source强制校验'),
    ('古典与现代脱节', '改"搜索"为"化用"，理解情感内核再匹配 + 白话解读'),
    ('海外IP被封', '切硅基流动国内代理 + Secrets管理Key + 超时异常处理'),
    ('HTML渲染乱码', 'f-string替代拼接 + html.escape()转义 + 独立CSS类'),
]
for i, (prob, sol) in enumerate(problems):
    y = 2.7 + i * 1.15
    rect(s, 1, y, 5.5, 0.95, LGRAY)
    tb(s, 1.2, y+0.08, 5, 0.35, prob, sz=15, b=True, c=GREEN_D)
    tb(s, 1.2, y+0.45, 5, 0.4, sol, sz=12, c=DARK)

# Right: reflections
tb(s, 7, 1.8, 5.5, 0.5, '核心反思', sz=20, b=True, c=GREEN_D)
rect(s, 7, 2.35, 5.5, 0.02, GREEN_L)
reflections = [
    'AI产品的价值不在于用了多少技术，\n而在于解决了什么真实问题。',
    '最难的不是调用API，而是想清楚\nAI应该做什么、不应该做什么。',
    '演示稳定性 > 功能完整性。\n简单不出错比复杂会翻车更有用。',
    '场景共鸣是最大的杠杆。\n"用古典回应现代情绪"，人人能懂。',
]
for i, ref in enumerate(reflections):
    y = 2.7 + i * 1.15
    rect(s, 7, y, 5.5, 0.95, GREEN_P)
    tbm(s, 7.2, y+0.1, 5, 0.75, ref.split('\n'), sz=13, c=DARK)

# ═══════════════════════════════════════
# SLIDE 6: THANK YOU
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, GREEN)

circle(s, -1, -1, 4, GREEN_D)
circle(s, 11, 5, 3.5, GREEN_D)

tb(s, 1, 2.0, 11, 1.2, '感谢聆听', sz=56, b=True, c=WHITE, a=PP_ALIGN.CENTER)
rect(s, 5, 3.3, 3.3, 0.05, WHITE)
tb(s, 1, 3.6, 11, 0.6, '拾句 · 一词入，千年应', sz=22, c=RGBColor(0xC8, 0xE6, 0xC9), a=PP_ALIGN.CENTER)

tb(s, 1, 4.8, 11, 0.5, '王  宸    2025212714', sz=18, c=WHITE, a=PP_ALIGN.CENTER)
tb(s, 1, 5.4, 11, 0.4, '北京邮电大学 · AI智能体开发', sz=14, c=RGBColor(0xA5, 0xD6, 0xA7), a=PP_ALIGN.CENTER)
tb(s, 1, 6.1, 11, 0.4,
    'ai-agent-work-as2jgjymz7emt2w5avocqv.streamlit.app',
    sz=12, c=RGBColor(0x81, 0xC7, 0x84), a=PP_ALIGN.CENTER)

prs.save(r'C:\Users\ZhiRan\Desktop\ppt\slides.pptx')
print('PPT DONE - 6 slides')
